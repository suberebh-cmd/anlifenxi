from __future__ import annotations

import generate_case_ppt as base
from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION, XL_LEGEND_POSITION
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


# Re-skin the existing helper module to fit the academic-pptx rules:
# white background, one font, restrained color usage.
base.BG = RGBColor(255, 255, 255)
base.INK = RGBColor(45, 45, 45)
base.SUBTLE = RGBColor(110, 110, 110)
base.ACCENT = RGBColor(31, 78, 121)
base.ACCENT_LIGHT = RGBColor(235, 243, 250)
base.ACCENT_WARM = RGBColor(46, 117, 182)
base.ACCENT_GOLD = RGBColor(255, 242, 204)
base.WHITE = RGBColor(255, 255, 255)
base.MUTED = RGBColor(204, 204, 204)
base.MUTED_DARK = RGBColor(140, 140, 140)
base.FONT_BODY = "Microsoft YaHei"
base.FONT_DISPLAY = "Microsoft YaHei"

FONT = base.FONT_BODY
BG = base.BG
INK = base.INK
SUBTLE = base.SUBTLE
ACCENT = base.ACCENT
ACCENT_LIGHT = base.ACCENT_LIGHT
ACCENT_WARM = base.ACCENT_WARM
ACCENT_GOLD = base.ACCENT_GOLD
WHITE = base.WHITE
MUTED = base.MUTED
MUTED_DARK = base.MUTED_DARK

set_slide_bg = base.set_slide_bg
add_rect = base.add_rect
add_textbox = base.add_textbox
add_paragraphs = base.add_paragraphs
add_bullets = base.add_bullets
add_footer = base.add_footer
add_slide_title = base.add_slide_title
add_metric_card = base.add_metric_card
style_chart = base.style_chart
add_column_chart = base.add_column_chart
add_line_chart = base.add_line_chart
add_bar_compare_chart = base.add_bar_compare_chart
add_single_bar_chart = base.add_single_bar_chart
add_question_box = base.add_question_box
add_step_card = base.add_step_card
read_csv_rows = base.read_csv_rows


def find_csv(prefix: str) -> Path:
    return next(p for p in Path("analysis_tables").iterdir() if p.name.startswith(prefix))


def add_divider_slide(prs: Presentation, section: str, title: str, subtitle: str, slide_no: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = ACCENT
    add_textbox(
        slide,
        Inches(0.9),
        Inches(1.2),
        Inches(2.8),
        Inches(0.45),
        section,
        font_size=14,
        color=RGBColor(190, 214, 237),
        bold=True,
    )
    add_textbox(
        slide,
        Inches(0.9),
        Inches(1.85),
        Inches(8.8),
        Inches(1.2),
        title,
        font_size=30,
        color=WHITE,
        bold=True,
    )
    add_textbox(
        slide,
        Inches(0.92),
        Inches(3.15),
        Inches(7.8),
        Inches(0.7),
        subtitle,
        font_size=17,
        color=RGBColor(220, 232, 245),
    )
    add_rect(slide, Inches(0.92), Inches(4.35), Inches(2.2), Inches(0.06), ACCENT_GOLD, ACCENT_GOLD)
    add_footer(slide, slide_no, "结构分节页")


def add_quote_band(slide, text: str) -> None:
    add_rect(slide, Inches(0.9), Inches(5.9), Inches(11.45), Inches(0.62), ACCENT_LIGHT, ACCENT_LIGHT, radius=True)
    add_textbox(
        slide,
        Inches(1.12),
        Inches(6.08),
        Inches(11.0),
        Inches(0.24),
        text,
        font_size=14,
        color=ACCENT,
        bold=True,
        align=PP_ALIGN.CENTER,
    )


def add_multi_line_chart(
    slide,
    left,
    top,
    width,
    height,
    categories,
    series_map: dict[str, list[float]],
    colors: list[RGBColor],
    legend_font_size: int = 10,
    show_labels: bool = False,
):
    data = CategoryChartData()
    data.categories = categories
    for name, values in series_map.items():
        data.add_series(name, values)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE_MARKERS, left, top, width, height, data
    ).chart
    style_chart(chart, legend=True)
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(legend_font_size)
    chart.legend.font.name = FONT
    for idx, series in enumerate(chart.series):
        color = colors[idx % len(colors)]
        series.format.line.color.rgb = color
        series.marker.style = 2
        series.marker.size = 6
        series.marker.format.fill.solid()
        series.marker.format.fill.fore_color.rgb = color
        if show_labels:
            series.has_data_labels = True
            series.data_labels.position = XL_LABEL_POSITION.ABOVE
            series.data_labels.font.size = Pt(8)
            series.data_labels.font.name = FONT
            series.data_labels.font.color.rgb = SUBTLE
    return chart


def add_single_series_bar_chart(
    slide,
    left,
    top,
    width,
    height,
    categories,
    values,
    series_name: str,
    color: RGBColor = ACCENT,
    highlight_indices: set[int] | None = None,
):
    data = CategoryChartData()
    data.categories = categories
    data.add_series(series_name, values)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, left, top, width, height, data
    ).chart
    style_chart(chart, legend=False)
    series = chart.series[0]
    series.has_data_labels = True
    series.data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
    series.data_labels.font.size = Pt(9)
    series.data_labels.font.name = FONT
    series.data_labels.font.color.rgb = SUBTLE
    highlight_indices = highlight_indices or set()
    for idx, point in enumerate(series.points):
        fill = point.format.fill
        fill.solid()
        fill.fore_color.rgb = color if idx in highlight_indices or not highlight_indices else MUTED_DARK
        point.format.line.color.rgb = fill.fore_color.rgb
    return chart


def add_three_point_summary(slide, bullets: list[str]) -> None:
    add_rect(slide, Inches(7.95), Inches(1.45), Inches(4.4), Inches(4.6), WHITE, MUTED, radius=True)
    add_bullets(
        slide,
        Inches(8.2),
        Inches(1.85),
        Inches(3.9),
        Inches(3.6),
        bullets,
        font_size=16,
    )


def build_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    annual_rows = read_csv_rows(find_csv("01_"))
    share_rows = read_csv_rows(find_csv("03_"))
    compare_rows = read_csv_rows(find_csv("04_"))
    t_rows = read_csv_rows(find_csv("06_"))
    copy_rows = read_csv_rows(find_csv("08_"))
    special_rows = read_csv_rows(find_csv("09_"))
    lag_rows = read_csv_rows(find_csv("10_"))

    years = [str(int(float(r["采购年"]))) for r in annual_rows]
    new_vols = [int(float(r["新增册数"])) for r in annual_rows]
    avg_prices = [float(r["均价"]) for r in annual_rows]
    median_prices = [float(r["中位价"]) for r in annual_rows]
    p90_prices = [float(r["P90价格"]) for r in annual_rows]
    one_year_ratio = [float(r["1年内%"]) for r in lag_rows]
    three_year_plus = [float(r["3年以上%"]) for r in lag_rows]
    single_copy_ratio = [float(r["单册品种占比%"]) for r in copy_rows]
    multi_copy_ratio = [float(r["2册及以上品种占比%"]) for r in copy_rows]

    compare_lookup = {r["一级类"]: r for r in compare_rows}
    compare_categories = [f"{c}\n{compare_lookup[c]['类别名称']}" for c in ["B", "C", "D", "I", "F", "T"]]
    compare_2021 = [float(compare_lookup[c]["2021占比%"]) for c in ["B", "C", "D", "I", "F", "T"]]
    compare_2025 = [float(compare_lookup[c]["2025占比%"]) for c in ["B", "C", "D", "I", "F", "T"]]

    share_map = {
        "B 哲学宗教": [float(r["B"]) for r in share_rows],
        "C 社科总论": [float(r["C"]) for r in share_rows],
        "D 政法": [float(r["D"]) for r in share_rows],
        "I 文学": [float(r["I"]) for r in share_rows],
        "F 经济": [float(r["F"]) for r in share_rows],
        "T 工业技术": [float(r["T"]) for r in share_rows],
    }

    t_focus_map = {
        "TP 自动化/计算机": [int(r["TP"]) for r in t_rows],
        "TN 无线电电子": [int(r["TN"]) for r in t_rows],
        "TS 轻工业": [int(r["TS"]) for r in t_rows],
        "TU 建筑科学": [int(r["TU"]) for r in t_rows],
    }

    location_names = [r["馆藏地"] for r in special_rows]
    location_counts = [int(float(r["记录数"])) for r in special_rows]
    location_avg_prices = [float(r["均价"]) for r in special_rows]

    ethnic_year_counts = [1919, 1504, 2247, 487, 951]
    ethnic_subject_labels = [
        "K 历史地理",
        "I 文学",
        "J 艺术",
        "Z 综合性图书",
        "G 文化教育",
        "F 经济",
        "D 政治法律",
        "H 语言文字",
        "C 社科总论",
        "R 医药卫生",
    ]
    ethnic_subject_counts = [2332, 1065, 839, 445, 428, 408, 395, 322, 256, 231]

    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = ACCENT
    add_textbox(slide, Inches(0.9), Inches(0.85), Inches(2.8), Inches(0.34), "academic-pptx 45分钟汇报版", 12, RGBColor(198, 219, 239), True)
    add_textbox(
        slide,
        Inches(0.9),
        Inches(1.45),
        Inches(8.8),
        Inches(1.5),
        "馆藏结构优化已经从“扩数量”\n变成“排优先级”",
        28,
        WHITE,
        True,
    )
    add_textbox(
        slide,
        Inches(0.92),
        Inches(3.15),
        Inches(8.4),
        Inches(0.9),
        "H大学图书馆馆藏结构优化教学型案例汇报\n一般学科资源与民族地区相关资源如何兼顾",
        17,
        RGBColor(220, 232, 245),
    )
    add_rect(slide, Inches(9.2), Inches(1.15), Inches(3.2), Inches(4.4), WHITE, WHITE, radius=True)
    add_paragraphs(
        slide,
        Inches(9.45),
        Inches(1.55),
        Inches(2.65),
        Inches(3.3),
        [
            ("这版结构按 45 分钟来做。", 16, ACCENT, True),
            ("我会先解释数据告诉了我们什么。", 16, INK, False),
            ("再把真正的决策难题抬出来。", 16, INK, False),
            ("最后回到教学型案例最重要的部分：不同方案各自会失去什么。", 16, INK, False),
        ],
    )
    add_textbox(slide, Inches(0.92), Inches(6.35), Inches(8.0), Inches(0.3), "汇报人：Codex 协作生成稿  |  场景：案例大赛 / 课程展示 / 馆务讨论", 11, RGBColor(198, 219, 239))
    add_footer(slide, 1, "版本：academic-pptx 快速长版")

    # Slide 2: Roadmap
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_slide_title(slide, "路线", "这 45 分钟，我想把问题讲成四个连续动作")
    roadmap = [
        ("01 看变化", "先看总量、价格、时效和复本结构，判断全馆资源配置到底变了什么。"),
        ("02 看迁移", "再看学科结构怎么移动，尤其是 B/C/D/I 上升和 T/F 回落意味着什么。"),
        ("03 看特色", "把民族地区相关资源单独拉出来，判断它们为什么不能继续淹没在平均数里。"),
        ("04 做选择", "最后回到案例核心：有限预算下，到底该用统一口径，还是分层治理。"),
    ]
    for idx, (label, body) in enumerate(roadmap):
        top = Inches(1.55 + idx * 1.15)
        add_rect(slide, Inches(0.95), top, Inches(1.2), Inches(0.72), ACCENT_LIGHT, ACCENT_LIGHT, radius=True)
        add_textbox(slide, Inches(1.2), top + Inches(0.18), Inches(0.7), Inches(0.2), label, 16, ACCENT, True, align=PP_ALIGN.CENTER)
        add_rect(slide, Inches(2.45), top, Inches(9.85), Inches(0.72), WHITE, MUTED, radius=True)
        add_textbox(slide, Inches(2.7), top + Inches(0.14), Inches(9.3), Inches(0.38), body, 16, INK)
    add_quote_band(slide, "如果只记住一句话，我希望是：这份案例真正讨论的不是“买多少”，而是“先保什么”。")
    add_footer(slide, 2, "内容结构依据：academic-pptx / content_guidelines.md")

    # Slide 3: Context
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_slide_title(slide, "情境", "馆藏优化已经从“增量管理”进入“优先级管理”")
    add_metric_card(slide, Inches(0.85), Inches(1.55), Inches(1.85), Inches(1.18), "142,847", "五年新增馆藏", "2021—2025")
    add_metric_card(slide, Inches(2.95), Inches(1.55), Inches(1.85), Inches(1.18), "99.04", "2025 平均单价", "元/册")
    add_metric_card(slide, Inches(5.05), Inches(1.55), Inches(1.85), Inches(1.18), "59.55%", "2025 一年内出版占比", "时效已修复")
    add_rect(slide, Inches(7.35), Inches(1.45), Inches(5.1), Inches(4.7), WHITE, MUTED, radius=True)
    add_bullets(
        slide,
        Inches(7.6),
        Inches(1.85),
        Inches(4.6),
        Inches(3.5),
        [
            "预算没有无限增长，空间也没有无限扩张。",
            "学校既要支持一般教学科研，也要回应民族地区人才培养、学科建设与文化传承任务。",
            "因此，馆藏优化不再是“尽量多买”，而是“必须排序”。",
        ],
        font_size=17,
    )
    add_rect(slide, Inches(0.95), Inches(3.15), Inches(5.95), Inches(2.05), ACCENT_LIGHT, ACCENT_LIGHT, radius=True)
    add_paragraphs(
        slide,
        Inches(1.2),
        Inches(3.45),
        Inches(5.4),
        Inches(1.4),
        [
            ("这个案例最有教学价值的地方在于：", 17, ACCENT, True),
            ("所有人都知道馆藏要优化，但一旦预算和评价同时受限，优化就会变成一道必须表态的选择题。", 18, INK, True),
        ],
    )
    add_footer(slide, 3, "数据来源：analysis_tables/01_年度总体统计.csv；案例分析初稿_馆藏结构优化_详析版.docx")

    # Slide 4: Research question
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_slide_title(slide, "核心问题", "本案例真正要问的是：有限预算下谁该先被保障")
    add_rect(slide, Inches(1.15), Inches(1.65), Inches(11.0), Inches(1.75), ACCENT_LIGHT, ACCENT, radius=True)
    add_textbox(
        slide,
        Inches(1.45),
        Inches(2.02),
        Inches(10.4),
        Inches(1.0),
        "当一般学科资源、课程复本需求、采购时效要求与民族地区相关资源建设同时存在时，\nH大学图书馆究竟应该继续使用统一评价口径，还是建立分层治理机制？",
        21,
        ACCENT,
        True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )
    add_bullets(
        slide,
        Inches(1.45),
        Inches(3.95),
        Inches(10.3),
        Inches(1.7),
        [
            "这不是抽象的价值讨论，而是预算分配、采购规则、馆藏评价和服务设计的组合决策。",
            "我后面所有图表都只服务这个问题，不追求把所有数据都讲完。",
        ],
        font_size=18,
    )
    add_footer(slide, 4, "结构要求：academic-pptx / 研究问题须在第2-3页前明确")

    # Slide 5: Data frame
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_slide_title(slide, "分析框架", "这份汇报只用三组证据来支撑判断")
    frame_items = [
        ("全馆运行", "新增册数、价格、时效、复本结构，用来判断馆藏运行状态是否稳定。"),
        ("结构变化", "一级类占比与 T 类二级拆解，用来判断资源配置方向有没有真正改变。"),
        ("特色资源", "民族地区相关馆藏规模、价格、空间功能与学科分布，用来判断其是否应被单列治理。"),
    ]
    for idx, (title, body) in enumerate(frame_items):
        left = Inches(0.95 + idx * 4.1)
        add_rect(slide, left, Inches(1.8), Inches(3.55), Inches(2.7), WHITE, MUTED, radius=True)
        add_textbox(slide, left + Inches(0.18), Inches(2.08), Inches(3.15), Inches(0.35), title, 18, ACCENT, True)
        add_textbox(slide, left + Inches(0.18), Inches(2.52), Inches(3.12), Inches(1.45), body, 16, INK)
    add_rect(slide, Inches(1.2), Inches(5.1), Inches(10.8), Inches(0.9), ACCENT_LIGHT, ACCENT_LIGHT, radius=True)
    add_textbox(
        slide,
        Inches(1.45),
        Inches(5.38),
        Inches(10.2),
        Inches(0.28),
        "判断标准也会跟着变化：常规新书更看覆盖与时效，特色资源更看使命承接与长期辨识度。",
        18,
        ACCENT,
        True,
        align=PP_ALIGN.CENTER,
    )
    add_footer(slide, 5, "数据来源：analysis_tables/01/03/04/08/09/10；案例分析初稿_民族文献聚焦版.docx")

    add_divider_slide(prs, "Part 1", "先看全馆运行面：总量没有停，治理压力却更高了", "这部分先回答一个前提问题：图书馆是不是在“失速”？", 6)

    # Slide 7: acquisitions
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_slide_title(slide, "总体变化", "五年新增馆藏有波动，但从未真正停下来")
    add_column_chart(slide, Inches(0.85), Inches(1.5), Inches(7.0), Inches(4.7), years, "新增册数", new_vols, ACCENT)
    add_three_point_summary(
        slide,
        [
            "2022 年是低点，但新增量仍有 25,630 册。",
            "2025 年回升到 29,979 册，说明全馆运行并未停摆。",
            "真正变化的不是“有没有买”，而是“往哪里买”。",
        ],
    )
    add_footer(slide, 7, "数据来源：analysis_tables/01_年度总体统计.csv")

    # Slide 8: price
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_slide_title(slide, "价格压力", "价格整体上行意味着“买错书”的代价变得更高")
    add_line_chart(slide, Inches(0.85), Inches(1.5), Inches(7.0), Inches(4.65), years, "平均单价（元）", avg_prices, ACCENT_WARM)
    add_metric_card(slide, Inches(8.05), Inches(1.55), Inches(1.95), Inches(1.05), "68→78", "中位价变化", "2021 到 2025")
    add_metric_card(slide, Inches(10.25), Inches(1.55), Inches(1.95), Inches(1.05), "119→129", "P90 价格变化", "高价书抬升")
    add_rect(slide, Inches(8.05), Inches(3.0), Inches(4.15), Inches(2.7), WHITE, MUTED, radius=True)
    add_bullets(
        slide,
        Inches(8.3),
        Inches(3.35),
        Inches(3.65),
        Inches(1.95),
        [
            "2022 与 2025 的均价都接近 100 元/册。",
            "价格上涨会放大资源配置错误的后果。",
            "所以后面的结构判断，不能再只看册数。",
        ],
        font_size=16,
    )
    add_footer(slide, 8, "数据来源：analysis_tables/01_年度总体统计.csv")

    # Slide 9: timeliness
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_slide_title(slide, "采购时效", "采购响应速度在 2023—2024 年明显变慢，2025 年才开始修复")
    add_multi_line_chart(
        slide,
        Inches(0.85),
        Inches(1.5),
        Inches(7.0),
        Inches(4.7),
        years,
        {"1年内出版占比%": one_year_ratio, "3年以上占比%": three_year_plus},
        [ACCENT_WARM, MUTED_DARK],
        legend_font_size=9,
        show_labels=True,
    )
    add_three_point_summary(
        slide,
        [
            "1 年内出版占比从 2021 年的 78.85% 下滑到 2023 年的 44.66%。",
            "2024 年 3 年以上旧书占比升到 21.35%，响应迟滞最明显。",
            "2025 年有所修复，但还没有回到 2021 年水平。",
        ],
    )
    add_footer(slide, 9, "数据来源：analysis_tables/10_出版时差分布.csv")

    # Slide 10: copy structure
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_slide_title(slide, "复本结构", "课程与服务型需求仍然很强，全馆并没有转向“只买研究书”")
    add_multi_line_chart(
        slide,
        Inches(0.85),
        Inches(1.5),
        Inches(7.0),
        Inches(4.7),
        years,
        {"单册品种占比%": single_copy_ratio, "2册及以上品种占比%": multi_copy_ratio},
        [ACCENT_WARM, ACCENT],
        legend_font_size=9,
        show_labels=True,
    )
    add_three_point_summary(
        slide,
        [
            "2 册及以上品种占比始终在 84% 以上，说明课程复本和普适性服务依旧占主流。",
            "2023 年单册品种占比抬升到 15.38%，像是一次结构试探。",
            "2025 年又回到 8.51%，表明常规教学支持不能被轻易压缩。",
        ],
    )
    add_footer(slide, 10, "数据来源：analysis_tables/08_复本结构_年度统计.csv")

    # Slide 11: interim summary
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_slide_title(slide, "阶段判断", "到这里我更愿意把问题概括成三句判断")
    add_rect(slide, Inches(0.95), Inches(1.7), Inches(11.45), Inches(4.25), WHITE, MUTED, radius=True)
    summary_text = [
        ("判断一", "总量没有停，所以馆藏并非“失血型”问题。"),
        ("判断二", "价格和时效压力同时存在，所以资源错误配置的成本正在上升。"),
        ("判断三", "复本结构没有明显瓦解，因此后面看到的结构迁移，更像优先级重排，而不是普通需求消失。"),
    ]
    for idx, (head, body) in enumerate(summary_text):
        top = Inches(2.05 + idx * 1.15)
        add_rect(slide, Inches(1.25), top, Inches(1.2), Inches(0.62), ACCENT, ACCENT, radius=True)
        add_textbox(slide, Inches(1.48), top + Inches(0.15), Inches(0.75), Inches(0.2), head, 15, WHITE, True, align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(2.7), top + Inches(0.1), Inches(8.9), Inches(0.36), body, 18, INK, True if idx == 2 else False)
    add_quote_band(slide, "换句话说，真正值得追问的已经不是“有没有买”，而是“哪些资源开始更重要”。")
    add_footer(slide, 11, "阶段小结：由 analysis_tables/01/08/10 归纳")

    add_divider_slide(prs, "Part 2", "再看结构迁移：资源排序到底朝哪个方向动了", "这一段把“结构变了”说具体，而不是停留在直觉上。", 12)

    # Slide 13: 2021 vs 2025
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_slide_title(slide, "结构对比", "2021 到 2025 年，资源配置更明显地转向人文社科议题")
    add_bar_compare_chart(
        slide,
        Inches(0.85),
        Inches(1.5),
        Inches(7.1),
        Inches(4.8),
        compare_categories,
        {"2021占比%": compare_2021, "2025占比%": compare_2025},
    )
    add_three_point_summary(
        slide,
        [
            "B 类从 6.49% 升到 12.94%，是最明显的上升项。",
            "C、D、I 也都上行，说明并非单一学科波动。",
            "F 与 T 的下降提示：原来较强的经管与技术板块正在让位。",
        ],
    )
    add_footer(slide, 13, "数据来源：analysis_tables/04_中图法一级类_2021_vs_2025.csv")

    # Slide 14: yearly trend
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_slide_title(slide, "连续趋势", "这不是某一年的偶然波动，而是连续五年的排序调整")
    add_multi_line_chart(
        slide,
        Inches(0.75),
        Inches(1.45),
        Inches(7.35),
        Inches(4.8),
        years,
        share_map,
        [ACCENT, ACCENT_WARM, RGBColor(112, 173, 71), RGBColor(91, 155, 213), RGBColor(165, 165, 165), RGBColor(192, 80, 77)],
        legend_font_size=8,
        show_labels=False,
    )
    add_three_point_summary(
        slide,
        [
            "B 类在 2025 年出现跳升，不像短期偶发采购。",
            "D 与 I 的占比长期维持高位，说明政法与文学持续受重视。",
            "F 与 T 不是某一年突然缩减，而是逐步承压。",
        ],
    )
    add_footer(slide, 14, "数据来源：analysis_tables/03_中图法一级类_年度占比.csv")

    # Slide 15: interpretive text
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_slide_title(slide, "结构解释", "上升的不只是“人文”，更是与学校使命更贴近的议题")
    add_rect(slide, Inches(0.95), Inches(1.55), Inches(5.7), Inches(4.8), WHITE, MUTED, radius=True)
    add_rect(slide, Inches(6.85), Inches(1.55), Inches(5.5), Inches(4.8), WHITE, MUTED, radius=True)
    add_textbox(slide, Inches(1.2), Inches(1.9), Inches(5.15), Inches(0.35), "从数据上看", 18, ACCENT, True)
    add_bullets(
        slide,
        Inches(1.2),
        Inches(2.35),
        Inches(4.95),
        Inches(3.35),
        [
            "B、C、D、I 都在上升，说明这轮变化更像价值主题和治理议题进入采购前列。",
            "K 类历史地理没有同步大涨，说明并非“凡是人文都增加”。",
            "这会让后面民族地区资源的判断更有解释力。",
        ],
        font_size=17,
    )
    add_textbox(slide, Inches(7.1), Inches(1.9), Inches(4.95), Inches(0.35), "从学校任务上看", 18, ACCENT_WARM, True)
    add_bullets(
        slide,
        Inches(7.1),
        Inches(2.35),
        Inches(4.8),
        Inches(3.35),
        [
            "民族地区人才培养、文化传承、共同体意识教育，本身就会把 B/C/D/I 推到更重要的位置。",
            "所以这不是“偏文科”，而是“任务导向更明显”。",
            "如果沿着这条线继续看，特色资源就不能再被视作边角补充。",
        ],
        font_size=17,
    )
    add_footer(slide, 15, "解释依据：analysis_tables/03/04；案例分析初稿_馆藏结构优化_详析版.docx")

    # Slide 16: T detailed compare
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_slide_title(slide, "T 类拆解", "T 类的回落主要发生在 TP 和 TN，而不是所有技术资源一起走弱")
    add_multi_line_chart(
        slide,
        Inches(0.85),
        Inches(1.45),
        Inches(7.1),
        Inches(4.8),
        years,
        t_focus_map,
        [ACCENT, ACCENT_WARM, RGBColor(112, 173, 71), RGBColor(91, 155, 213)],
        legend_font_size=8,
        show_labels=False,
    )
    add_three_point_summary(
        slide,
        [
            "TP 从 2103 册降到 1115 册，回落最明显。",
            "TN 也从 587 册降到 212 册，技术硬件相关支持同步承压。",
            "TS、TU 的波动较小，说明并不是所有技术类都在退场。",
        ],
    )
    add_footer(slide, 16, "数据来源：analysis_tables/06_T类二级分类_年度册数.csv")

    # Slide 17: T interpretation
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_slide_title(slide, "再解释一步", "技术资源并没有失去必要性，但它们不再自动拥有优先级")
    add_rect(slide, Inches(0.95), Inches(1.7), Inches(11.45), Inches(4.25), WHITE, MUTED, radius=True)
    bullets = [
        "第一，TP 仍然是 T 类主体，说明数字化、计算机类资源仍有刚性需求。",
        "第二，回落更多像是相对优先级下降，而不是绝对需求消失。",
        "第三，问题因此变得更难：我们不能简单说“多买人文、少买技术”，而要解释哪些任务必须被优先保障。",
    ]
    add_bullets(slide, Inches(1.35), Inches(2.1), Inches(10.5), Inches(2.4), bullets, font_size=19)
    add_quote_band(slide, "这一步很关键，因为它让案例避免滑向“文科替代工科”这种过于粗糙的叙事。")
    add_footer(slide, 17, "解释依据：analysis_tables/06_T类二级分类_年度册数.csv")

    # Slide 18: section summary
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_slide_title(slide, "阶段判断", "到这里，我会把结构变化概括成一个更具体的结论")
    add_rect(slide, Inches(1.0), Inches(1.8), Inches(11.2), Inches(1.3), ACCENT_LIGHT, ACCENT_LIGHT, radius=True)
    add_textbox(
        slide,
        Inches(1.3),
        Inches(2.15),
        Inches(10.6),
        Inches(0.65),
        "H大学近五年的资源配置，并不是在简单压缩某些学科，而是在把采购优先级从“默认平均分配”转向“更贴近学校任务的板块”。",
        21,
        ACCENT,
        True,
        align=PP_ALIGN.CENTER,
    )
    add_bullets(
        slide,
        Inches(1.25),
        Inches(3.65),
        Inches(10.4),
        Inches(1.9),
        [
            "这为民族地区相关资源的独立讨论提供了前提。",
            "如果全馆都在重排优先级，那么特色资源就更应该被问清楚：它们到底只是补充，还是应该被保障。",
        ],
        font_size=18,
    )
    add_footer(slide, 18, "阶段小结：由 analysis_tables/03/04/06 归纳")

    add_divider_slide(prs, "Part 3", "把民族地区相关资源单独拉出来，它们的意义会更清楚", "这部分不再问“有没有”，而是问“为什么不能继续混算”。", 19)

    # Slide 20: ethnic scale
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_slide_title(slide, "规模判断", "民族地区相关资源只占 4.98%，却最能代表学校是谁")
    add_metric_card(slide, Inches(0.85), Inches(1.5), Inches(1.85), Inches(1.18), "9015", "相关馆藏总记录", "专题统计")
    add_metric_card(slide, Inches(2.95), Inches(1.5), Inches(1.85), Inches(1.18), "7108", "五年新增", "2021—2025")
    add_metric_card(slide, Inches(5.05), Inches(1.5), Inches(1.85), Inches(1.18), "4.98%", "占同期新增比重", "并不高")
    add_column_chart(slide, Inches(0.85), Inches(3.0), Inches(6.95), Inches(3.0), years, "年度新增册数", ethnic_year_counts, ACCENT)
    add_three_point_summary(
        slide,
        [
            "规模不大，所以它们最容易在统一平均口径里被忽略。",
            "但它们高度集中承接学校特色任务，战略分量远高于册数占比。",
            "这正是“单列治理”比“继续混算”更有必要的原因。",
        ],
    )
    add_footer(slide, 20, "数据来源：案例分析初稿_民族文献聚焦版.docx；analysis_tables/01_年度总体统计.csv")

    # Slide 21: ethnic volatility
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_slide_title(slide, "年度波动", "民族地区相关资源在 2024 年明显收缩，说明它们最容易在统一预算里被让位")
    add_column_chart(slide, Inches(0.85), Inches(1.5), Inches(7.0), Inches(4.7), years, "年度新增册数", ethnic_year_counts, ACCENT_WARM)
    add_three_point_summary(
        slide,
        [
            "2023 年达到 2247 册后，2024 年骤降到 487 册。",
            "2025 年回升到 951 册，但仍显著低于 2023 年。",
            "这类大起大落，恰恰说明特色资源最怕被放进统一削减逻辑。",
        ],
    )
    add_footer(slide, 21, "数据来源：案例分析初稿_民族文献聚焦版.docx")

    # Slide 22: three spaces
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_slide_title(slide, "空间分工", "L203、L201 与 L102 的功能并不相同，不能再用同一把尺子衡量")
    cards = [
        ("L203 民族文献借阅", "更偏教学支持与可借服务。记录数 3888，均价 78.54 元。"),
        ("L201 民族文献阅览", "更偏专题研究与现场查阅。记录数 3804，均价 221.76 元。"),
        ("L102 民大文库", "更偏收藏与学校辨识度表达。记录数 1323，均价 262.04 元。"),
    ]
    for idx, (title, body) in enumerate(cards):
        left = Inches(0.95 + idx * 4.05)
        add_rect(slide, left, Inches(1.9), Inches(3.6), Inches(3.0), WHITE, MUTED, radius=True)
        add_textbox(slide, left + Inches(0.18), Inches(2.18), Inches(2.7), Inches(0.38), title, 17, ACCENT, True)
        add_textbox(slide, left + Inches(0.18), Inches(2.72), Inches(2.65), Inches(1.35), body, 16, INK)
    add_quote_band(slide, "借阅、阅览、收藏承担的是三种不同任务，所以它们不该继续挤在同一张成绩单上。")
    add_footer(slide, 22, "数据来源：analysis_tables/09_特色馆藏地点统计.csv")

    # Slide 23: location counts
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_slide_title(slide, "空间规模", "特色馆藏空间里，民族文献板块不是最大的，但它们的功能最专门")
    add_single_series_bar_chart(
        slide,
        Inches(0.85),
        Inches(1.45),
        Inches(7.0),
        Inches(4.9),
        location_names,
        location_counts,
        "记录数",
        color=ACCENT,
        highlight_indices={0, 1, 4},
    )
    add_three_point_summary(
        slide,
        [
            "多卷工具阅览与捐赠书空间规模更大，但承担的是另一类基础保障任务。",
            "民族文献借阅、阅览与民大文库数量有限，却承担更清晰的学校特色表达功能。",
            "所以，单看空间规模会系统性低估这部分资源的重要性。",
        ],
    )
    add_footer(slide, 23, "数据来源：analysis_tables/09_特色馆藏地点统计.csv")

    # Slide 24: location prices
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_slide_title(slide, "价格结构", "研究型与收藏型民族文献明显更贵，这决定了它们不能照搬常规新书规则")
    add_single_series_bar_chart(
        slide,
        Inches(0.85),
        Inches(1.45),
        Inches(7.0),
        Inches(4.9),
        location_names,
        location_avg_prices,
        "平均单价",
        color=ACCENT_WARM,
        highlight_indices={0, 1, 4},
    )
    add_three_point_summary(
        slide,
        [
            "L201 与 L102 的均价都超过 220 元/册，远高于全馆年度均价。",
            "这往往对应专题资料、影印文献和收藏型资源，而不是普通新书。",
            "若仍用统一成本标准评价，它们会天然处于劣势。",
        ],
    )
    add_footer(slide, 24, "数据来源：analysis_tables/09_特色馆藏地点统计.csv；analysis_tables/01_年度总体统计.csv")

    # Slide 25: subject composition
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_slide_title(slide, "学科面貌", "民族地区相关资源并不是单一 K 类，而是一组跨学科组合")
    add_single_series_bar_chart(
        slide,
        Inches(0.78),
        Inches(1.42),
        Inches(7.15),
        Inches(5.0),
        ethnic_subject_labels,
        ethnic_subject_counts,
        "记录数",
        color=ACCENT,
        highlight_indices={0, 1, 2, 3},
    )
    add_three_point_summary(
        slide,
        [
            "K 类历史地理是核心，但 I、J、C、D、G 都有明显分量。",
            "这意味着民族地区资源建设并不只是“补地方文献”。",
            "它更像一组围绕学校使命组织起来的跨学科资源包。",
        ],
    )
    add_footer(slide, 25, "数据来源：案例分析初稿_民族文献聚焦版.docx")

    # Slide 26: mission
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_slide_title(slide, "使命承接", "这些资源之所以重要，是因为它们承接的是学校的特殊任务")
    mission_items = [
        ("人才培养", "民族地区课程、专题阅读和学生学习支持需要稳定资源供给。"),
        ("学科研究", "民族学、历史、文化、治理等方向的研究深度离不开专题资料。"),
        ("文化传承", "影印文献、民间文学与地方文化资源具有长期保存价值。"),
        ("学校辨识度", "馆藏是办学定位最可见的物质证据之一。"),
    ]
    for idx, (title, body) in enumerate(mission_items):
        left = Inches(0.95 + (idx % 2) * 5.95)
        top = Inches(1.8 + (idx // 2) * 1.8)
        add_rect(slide, left, top, Inches(5.45), Inches(1.3), WHITE, MUTED, radius=True)
        add_textbox(slide, left + Inches(0.18), top + Inches(0.18), Inches(1.4), Inches(0.25), title, 17, ACCENT, True)
        add_textbox(slide, left + Inches(1.78), top + Inches(0.16), Inches(3.2), Inches(0.72), body, 15, INK)
    add_quote_band(slide, "一般资源回答“课堂上要用什么”，特色资源回答“这所学校到底是谁”。")
    add_footer(slide, 26, "解释依据：案例分析初稿_馆藏结构优化_详析版.docx；案例分析初稿_民族文献聚焦版.docx")

    add_divider_slide(prs, "Part 4", "最后回到决策：管理上到底该怎么选", "这部分不追求标准答案，而是把各个选项的代价讲清楚。", 27)

    # Slide 28: decision point
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_slide_title(slide, "决策时点", "馆藏结构优化真正卡住的，是四个必须同时回答的问题")
    questions = [
        ("预算要不要单列", "民族地区相关资源是否需要独立保障额度。"),
        ("评价要不要单列", "是否继续与常规新书一起计算平均数、时效与绩效。"),
        ("空间要不要分工", "借阅、阅览、收藏三类空间是否要承担不同目标。"),
        ("服务要不要转化", "专题资源是否要和课程、导读、展陈、数字化联动。"),
    ]
    for idx, (title, body) in enumerate(questions):
        left = Inches(0.95 + (idx % 2) * 5.95)
        top = Inches(1.75 + (idx // 2) * 1.85)
        add_question_box(slide, left, top, Inches(5.45), Inches(1.45), title, body)
    add_quote_band(slide, "教学型案例的张力，就在于每个问题都没有零成本答案。")
    add_footer(slide, 28, "决策议题归纳：案例分析初稿_馆藏结构优化_详析版.docx")

    # Slide 29: option A
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_slide_title(slide, "方案 A", "继续使用统一口径，管理会更省事，但学校辨识度会继续被稀释")
    add_rect(slide, Inches(0.95), Inches(1.7), Inches(5.5), Inches(4.4), WHITE, MUTED, radius=True)
    add_textbox(slide, Inches(1.2), Inches(2.0), Inches(1.0), Inches(0.25), "优势", 18, ACCENT, True)
    add_bullets(
        slide,
        Inches(1.2),
        Inches(2.35),
        Inches(4.7),
        Inches(2.8),
        [
            "统计规则简单，部门协同成本低。",
            "预算审批与绩效考核都更容易解释。",
            "不会打破既有采购流程。",
        ],
        font_size=17,
    )
    add_rect(slide, Inches(6.85), Inches(1.7), Inches(5.5), Inches(4.4), WHITE, MUTED, radius=True)
    add_textbox(slide, Inches(7.1), Inches(2.0), Inches(1.0), Inches(0.25), "代价", 18, ACCENT_WARM, True)
    add_bullets(
        slide,
        Inches(7.1),
        Inches(2.35),
        Inches(4.7),
        Inches(2.8),
        [
            "特色资源会继续显得“贵、慢、少”。",
            "一遇到预算收紧，它们会最先被让位。",
            "学校特色难以稳定沉淀成馆藏优势。",
        ],
        font_size=17,
    )
    add_footer(slide, 29, "方案分析：基于本案数据与教学讨论归纳")

    # Slide 30: option B
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_slide_title(slide, "方案 B", "把民族地区相关资源单列出来，会增加治理成本，却能换来更清晰的办学表达")
    add_rect(slide, Inches(0.95), Inches(1.7), Inches(5.5), Inches(4.4), WHITE, MUTED, radius=True)
    add_textbox(slide, Inches(1.2), Inches(2.0), Inches(1.0), Inches(0.25), "优势", 18, ACCENT, True)
    add_bullets(
        slide,
        Inches(1.2),
        Inches(2.35),
        Inches(4.7),
        Inches(2.8),
        [
            "能用更适合特色资源的标准来评价其价值。",
            "预算波动时更容易守住学校辨识度资源。",
            "能倒逼空间分工和服务转化真正发生。",
        ],
        font_size=17,
    )
    add_rect(slide, Inches(6.85), Inches(1.7), Inches(5.5), Inches(4.4), WHITE, MUTED, radius=True)
    add_textbox(slide, Inches(7.1), Inches(2.0), Inches(1.0), Inches(0.25), "代价", 18, ACCENT_WARM, True)
    add_bullets(
        slide,
        Inches(7.1),
        Inches(2.35),
        Inches(4.7),
        Inches(2.8),
        [
            "统计、论证和审批流程都会更复杂。",
            "需要更细的分类口径与责任分工。",
            "如果只单列不转服务，仍可能停留在“被看见但没被用好”。",
        ],
        font_size=17,
    )
    add_footer(slide, 30, "方案分析：基于本案数据与教学讨论归纳")

    # Slide 31: recommendation
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_slide_title(slide, "我的倾向", "我更倾向于“分层治理”，而不是继续追求一个总分")
    add_rect(slide, Inches(1.0), Inches(1.8), Inches(11.2), Inches(1.2), ACCENT_LIGHT, ACCENT_LIGHT, radius=True)
    add_textbox(
        slide,
        Inches(1.2),
        Inches(2.15),
        Inches(10.8),
        Inches(0.45),
        "常规新书看覆盖、时效与复本；特色资源看使命承接、长期积累与服务转化。",
        22,
        ACCENT,
        True,
        align=PP_ALIGN.CENTER,
    )
    add_bullets(
        slide,
        Inches(1.2),
        Inches(3.45),
        Inches(10.5),
        Inches(2.1),
        [
            "这不是给特色资源“特殊照顾”，而是承认不同资源承担的任务本来就不同。",
            "只要继续用统一口径，民族地区相关资源就会不断被误判为“低效率”。",
            "分层治理的前提不是多花钱，而是先把规则讲对。",
        ],
        font_size=18,
    )
    add_footer(slide, 31, "建议方案：基于本案数据与教学判断提出")

    # Slide 32: implementation
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_slide_title(slide, "实施路径", "如果要落地，我会先做五个动作，而不是一次性推翻现有体系")
    steps = [
        "先分口径：常规新增、特色馆藏、工具书、捐赠资源分开统计。",
        "再定清单：把民族地区相关资源的范围和馆藏地责任界定清楚。",
        "再设指标：特色资源单列覆盖度、连续性、使用转化与展示转化指标。",
        "再做联动：让课程、导读、展陈、数字化一起承接这些资源。",
        "最后复盘：每年不只看买了多少，更看学校特色有没有被真正看见和用起来。",
    ]
    positions = [
        (Inches(0.95), Inches(1.8)),
        (Inches(6.8), Inches(1.8)),
        (Inches(0.95), Inches(3.15)),
        (Inches(6.8), Inches(3.15)),
        (Inches(3.88), Inches(4.5)),
    ]
    for idx, (pos, text) in enumerate(zip(positions, steps), start=1):
        add_step_card(slide, pos[0], pos[1], Inches(5.45), Inches(1.0), str(idx), text)
    add_footer(slide, 32, "实施建议：基于本案教学型分析提出")

    # Slide 33: teaching discussion
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_slide_title(slide, "课堂讨论", "如果你是资源建设部负责人，你会怎样向馆长解释这 4.98%")
    add_rect(slide, Inches(0.95), Inches(1.7), Inches(11.35), Inches(4.4), WHITE, MUTED, radius=True)
    questions = [
        "1. 在预算只小幅增长时，你会优先守住常规时效，还是优先守住民族地区相关资源的连续建设？",
        "2. 特色资源单列后，最合适的绩效指标应该是什么，哪些指标反而不该继续沿用？",
        "3. 如果领导只看“平均单价”和“流通量”，你会怎么解释高价专题资源的合理性？",
    ]
    add_paragraphs(
        slide,
        Inches(1.25),
        Inches(2.08),
        Inches(10.7),
        Inches(2.85),
        [(q, 20, INK, True if idx == 0 else False) for idx, q in enumerate(questions)],
    )
    add_quote_band(slide, "我会把这份案例停在选择题上，因为教学型案例的价值本来就不在“唯一答案”。")
    add_footer(slide, 33, "用途：课堂讨论 / 案例大赛展示 / 馆务研讨")

    # Slide 34: conclusions
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_slide_title(slide, "结论", "我的结论是：H大学真正要优化的，不是总量，而是优先级")
    add_rect(slide, Inches(0.95), Inches(1.7), Inches(11.35), Inches(4.2), WHITE, MUTED, radius=True)
    conclusion_bullets = [
        "五年馆藏总体仍在运行，但价格与时效压力让“买对什么”比“买多少”更重要。",
        "资源结构已从平均分配转向任务导向，B/C/D/I 上升、F/T 承压就是这个信号。",
        "民族地区相关资源规模不大，却承担学校辨识度和特殊任务，因此不宜继续与常规新书混算。",
        "我更支持“分层治理 + 服务转化”的方案，让特色资源既被看见，也被真正用起来。",
    ]
    add_bullets(slide, Inches(1.2), Inches(2.15), Inches(10.75), Inches(2.95), conclusion_bullets, font_size=20)
    add_textbox(slide, Inches(1.25), Inches(5.35), Inches(3.3), Inches(0.22), "Q&A 停留页", 14, ACCENT_WARM, True)
    add_textbox(slide, Inches(8.7), Inches(5.33), Inches(2.9), Inches(0.22), "contact: academic-pptx working draft", 12, SUBTLE, align=PP_ALIGN.RIGHT)
    add_footer(slide, 34, "问答停留页")

    # Slide 35: references
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_slide_title(slide, "参考来源", "本次汇报使用的数据与文本来源如下")
    refs = [
        "1. analysis_tables/01_年度总体统计.csv",
        "2. analysis_tables/03_中图法一级类_年度占比.csv",
        "3. analysis_tables/04_中图法一级类_2021_vs_2025.csv",
        "4. analysis_tables/06_T类二级分类_年度册数.csv",
        "5. analysis_tables/08_复本结构_年度统计.csv",
        "6. analysis_tables/09_特色馆藏地点统计.csv",
        "7. analysis_tables/10_出版时差分布.csv",
        "8. 案例分析初稿_馆藏结构优化_详析版.docx",
        "9. 案例分析初稿_民族文献聚焦版.docx",
        "10. academic-pptx skill: SKILL.md / content_guidelines.md / slide_patterns.md",
    ]
    add_bullets(slide, Inches(1.15), Inches(1.8), Inches(10.9), Inches(4.7), refs, font_size=17)
    add_footer(slide, 35, "References")

    add_divider_slide(prs, "Appendix", "附录页保留给追问数据与补充图表", "正式汇报可按现场时间决定是否展开。", 36)

    # Slide 37: appendix annual table
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_slide_title(slide, "附录 A", "全馆年度运行数据把“规模、价格、时效”同时放在一起看更清楚")
    headers = ["年份", "新增册数", "均价", "中位价", "P90价", "1年内占比"]
    x_positions = [0.95, 2.1, 3.75, 5.15, 6.45, 7.75]
    widths = [0.8, 1.35, 1.05, 1.0, 1.0, 1.4]
    for x, w, header in zip(x_positions, widths, headers):
        add_rect(slide, Inches(x), Inches(1.75), Inches(w), Inches(0.5), ACCENT_LIGHT, ACCENT_LIGHT, radius=True)
        add_textbox(slide, Inches(x + 0.05), Inches(1.92), Inches(w - 0.1), Inches(0.18), header, 13, ACCENT, True, align=PP_ALIGN.CENTER)
    for idx, year in enumerate(years):
        top = Inches(2.35 + idx * 0.68)
        values = [year, f"{new_vols[idx]}", f"{avg_prices[idx]:.2f}", f"{median_prices[idx]:.1f}", f"{p90_prices[idx]:.1f}", f"{one_year_ratio[idx]:.2f}%"]
        for x, w, value in zip(x_positions, widths, values):
            add_rect(slide, Inches(x), top, Inches(w), Inches(0.48), WHITE, MUTED, radius=True)
            add_textbox(slide, Inches(x + 0.04), top + Inches(0.13), Inches(w - 0.08), Inches(0.18), value, 13, INK, False, align=PP_ALIGN.CENTER)
    add_footer(slide, 37, "数据来源：analysis_tables/01_年度总体统计.csv；analysis_tables/10_出版时差分布.csv")

    # Slide 38: appendix raw view
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_slide_title(slide, "附录 B", "特色馆藏空间的原始记录数与均价放在一起看")
    add_single_series_bar_chart(
        slide,
        Inches(0.75),
        Inches(1.45),
        Inches(5.9),
        Inches(4.8),
        location_names,
        location_counts,
        "记录数",
        color=ACCENT,
        highlight_indices={0, 1, 4},
    )
    add_single_series_bar_chart(
        slide,
        Inches(6.85),
        Inches(1.45),
        Inches(5.9),
        Inches(4.8),
        location_names,
        location_avg_prices,
        "平均单价",
        color=ACCENT_WARM,
        highlight_indices={0, 1, 4},
    )
    add_footer(slide, 38, "数据来源：analysis_tables/09_特色馆藏地点统计.csv")

    return prs


def main() -> None:
    prs = build_presentation()
    out = Path("馆藏结构优化_教学案例汇报_academic-pptx_45分钟版.pptx")
    prs.save(str(out))
    print(out.name)
    print(len(prs.slides))


if __name__ == "__main__":
    main()
