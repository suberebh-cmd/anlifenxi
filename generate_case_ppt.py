from __future__ import annotations

import csv
from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


FONT_BODY = "Microsoft YaHei"
FONT_DISPLAY = "Microsoft YaHei"

BG = RGBColor(248, 245, 239)
INK = RGBColor(30, 43, 56)
SUBTLE = RGBColor(111, 122, 133)
ACCENT = RGBColor(26, 66, 102)
ACCENT_LIGHT = RGBColor(216, 228, 235)
ACCENT_WARM = RGBColor(169, 70, 57)
ACCENT_GOLD = RGBColor(185, 140, 58)
WHITE = RGBColor(255, 255, 255)
MUTED = RGBColor(214, 212, 206)
MUTED_DARK = RGBColor(162, 159, 153)


def read_csv_rows(path: Path) -> list[dict[str, str]]:
    with path.open("r", encoding="utf-8-sig", newline="") as fh:
        return list(csv.DictReader(fh))


def set_slide_bg(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_rect(slide, left, top, width, height, fill_color, line_color=None, radius=False):
    shape_type = (
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE
        if radius
        else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    )
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = fill_color
    shape.line.color.rgb = line_color or fill_color
    return shape


def add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text="",
    font_size=20,
    color=INK,
    bold=False,
    font_name=FONT_BODY,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.name = font_name
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_paragraphs(
    slide,
    left,
    top,
    width,
    height,
    paragraphs: list[tuple[str, int, RGBColor, bool]],
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = valign
    tf.clear()
    for idx, (text, font_size, color, bold) in enumerate(paragraphs):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(4)
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.name = FONT_BODY
        run.font.bold = bold
        run.font.color.rgb = color
    return box


def add_bullets(slide, left, top, width, height, bullets: list[str], font_size=19):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for idx, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.space_after = Pt(8)
        p.font.size = Pt(font_size)
        p.font.name = FONT_BODY
        p.font.color.rgb = INK
        p.bullet = True
    return box


def add_footer(slide, slide_no: int, source_text: str = "数据来源：五年馆藏.xlsx / analysis_tables") -> None:
    add_rect(slide, Inches(0.55), Inches(7.05), Inches(12.25), Inches(0.02), MUTED)
    add_textbox(
        slide,
        Inches(0.58),
        Inches(7.1),
        Inches(5.8),
        Inches(0.22),
        source_text,
        font_size=9,
        color=SUBTLE,
    )
    add_textbox(
        slide,
        Inches(12.0),
        Inches(7.08),
        Inches(0.7),
        Inches(0.25),
        f"{slide_no:02d}",
        font_size=10,
        color=SUBTLE,
        bold=True,
        align=PP_ALIGN.RIGHT,
    )


def add_slide_title(slide, eyebrow: str, title: str, subtitle: str | None = None) -> None:
    add_textbox(
        slide,
        Inches(0.75),
        Inches(0.45),
        Inches(2.3),
        Inches(0.28),
        eyebrow,
        font_size=11,
        color=ACCENT_WARM,
        bold=True,
    )
    add_textbox(
        slide,
        Inches(0.72),
        Inches(0.75),
        Inches(7.8),
        Inches(0.75),
        title,
        font_size=26,
        color=INK,
        bold=True,
        font_name=FONT_DISPLAY,
    )
    if subtitle:
        add_textbox(
            slide,
            Inches(0.75),
            Inches(1.45),
            Inches(8.8),
            Inches(0.45),
            subtitle,
            font_size=13,
            color=SUBTLE,
        )


def add_metric_card(slide, left, top, width, height, value: str, label: str, note: str = "", fill=WHITE):
    add_rect(slide, left, top, width, height, fill, MUTED, radius=True)
    add_textbox(slide, left + Inches(0.14), top + Inches(0.08), width - Inches(0.28), Inches(0.4), value, 24, ACCENT, True)
    add_textbox(slide, left + Inches(0.14), top + Inches(0.48), width - Inches(0.28), Inches(0.36), label, 11, INK, True)
    if note:
        add_textbox(slide, left + Inches(0.14), top + Inches(0.82), width - Inches(0.28), Inches(0.28), note, 9, SUBTLE)


def style_chart(chart, legend=True):
    chart.chart_style = 2
    chart.has_legend = legend
    if legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(10)
        chart.legend.font.name = FONT_BODY
    chart.has_title = False

    if chart.category_axis:
        chart.category_axis.tick_labels.font.size = Pt(10)
        chart.category_axis.tick_labels.font.name = FONT_BODY
        chart.category_axis.tick_labels.font.color.rgb = SUBTLE

    if chart.value_axis:
        chart.value_axis.tick_labels.font.size = Pt(10)
        chart.value_axis.tick_labels.font.name = FONT_BODY
        chart.value_axis.tick_labels.font.color.rgb = SUBTLE
        chart.value_axis.has_major_gridlines = True
        chart.value_axis.major_gridlines.format.line.color.rgb = MUTED


def add_column_chart(slide, left, top, width, height, categories, series_name, values, color):
    data = CategoryChartData()
    data.categories = categories
    data.add_series(series_name, values)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, left, top, width, height, data
    ).chart
    style_chart(chart, legend=False)
    series = chart.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = color
    series.format.line.color.rgb = color
    series.has_data_labels = True
    series.data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
    series.data_labels.font.size = Pt(9)
    series.data_labels.font.name = FONT_BODY
    series.data_labels.font.color.rgb = SUBTLE
    return chart


def add_line_chart(slide, left, top, width, height, categories, series_name, values, color):
    data = CategoryChartData()
    data.categories = categories
    data.add_series(series_name, values)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE_MARKERS, left, top, width, height, data
    ).chart
    style_chart(chart, legend=False)
    series = chart.series[0]
    series.format.line.color.rgb = color
    series.marker.style = 2
    series.marker.size = 7
    series.marker.format.fill.solid()
    series.marker.format.fill.fore_color.rgb = color
    series.has_data_labels = True
    series.data_labels.position = XL_LABEL_POSITION.ABOVE
    series.data_labels.font.size = Pt(9)
    series.data_labels.font.name = FONT_BODY
    series.data_labels.font.color.rgb = SUBTLE
    return chart


def add_bar_compare_chart(slide, left, top, width, height, categories, series_map):
    data = CategoryChartData()
    data.categories = categories
    for name, values in series_map.items():
        data.add_series(name, values)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, left, top, width, height, data
    ).chart
    style_chart(chart, legend=True)
    colors = [ACCENT, ACCENT_WARM]
    for idx, series in enumerate(chart.series):
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = colors[idx]
        series.format.line.color.rgb = colors[idx]
    return chart


def add_single_bar_chart(slide, left, top, width, height, categories, values, highlight_indices):
    data = CategoryChartData()
    data.categories = categories
    data.add_series("记录数", values)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, left, top, width, height, data
    ).chart
    style_chart(chart, legend=False)
    series = chart.series[0]
    for idx, point in enumerate(series.points):
        fill = point.format.fill
        fill.solid()
        fill.fore_color.rgb = ACCENT if idx in highlight_indices else MUTED_DARK
    return chart


def add_question_box(slide, left, top, width, height, title: str, body: str):
    add_rect(slide, left, top, width, height, WHITE, MUTED, radius=True)
    add_textbox(slide, left + Inches(0.15), top + Inches(0.12), width - Inches(0.3), Inches(0.32), title, 13, ACCENT_WARM, True)
    add_textbox(slide, left + Inches(0.15), top + Inches(0.5), width - Inches(0.3), height - Inches(0.65), body, 13, INK)


def add_step_card(slide, left, top, width, height, num: str, text: str):
    add_rect(slide, left, top, width, height, WHITE, MUTED, radius=True)
    add_rect(slide, left + Inches(0.12), top + Inches(0.12), Inches(0.38), Inches(0.38), ACCENT, ACCENT, radius=True)
    add_textbox(
        slide,
        left + Inches(0.12),
        top + Inches(0.11),
        Inches(0.38),
        Inches(0.38),
        num,
        13,
        WHITE,
        True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )
    add_textbox(slide, left + Inches(0.62), top + Inches(0.1), width - Inches(0.74), height - Inches(0.2), text, 14, INK)


def build_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    base = Path("analysis_tables")
    annual_rows = read_csv_rows(base / "01_年度总体统计.csv")
    compare_rows = read_csv_rows(base / "04_中图法一级类_2021_vs_2025.csv")
    special_rows = read_csv_rows(base / "09_特色馆藏地点统计.csv")
    lag_rows = read_csv_rows(base / "10_出版时差分布.csv")

    years = [str(int(float(r["采购年"]))) for r in annual_rows]
    new_vols = [int(float(r["新增册数"])) for r in annual_rows]
    one_year_ratio = [float(r["一年内出版占比"]) for r in annual_rows]

    selected_cats = ["B", "C", "D", "I", "F", "T"]
    cat_lookup = {r["一级类"]: r for r in compare_rows}
    compare_categories = [f"{c}\n{cat_lookup[c]['类别名称']}" for c in selected_cats]
    compare_2021 = [float(cat_lookup[c]["2021占比%"]) for c in selected_cats]
    compare_2025 = [float(cat_lookup[c]["2025占比%"]) for c in selected_cats]

    location_names = [r["馆藏地"] for r in special_rows]
    location_counts = [int(float(r["记录数"])) for r in special_rows]

    # Slide 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_rect(slide, Inches(0.0), Inches(0.0), Inches(0.38), Inches(7.5), ACCENT)
    add_rect(slide, Inches(0.75), Inches(0.6), Inches(2.2), Inches(0.34), ACCENT_LIGHT, ACCENT_LIGHT, radius=True)
    add_textbox(slide, Inches(0.9), Inches(0.66), Inches(2.0), Inches(0.2), "教学型案例汇报初稿", 11, ACCENT, True)
    add_textbox(
        slide,
        Inches(0.78),
        Inches(1.2),
        Inches(8.0),
        Inches(1.5),
        "真正难的不是买多少书，\n而是先为谁留位置",
        30,
        INK,
        True,
        font_name=FONT_DISPLAY,
    )
    add_textbox(
        slide,
        Inches(0.82),
        Inches(2.9),
        Inches(7.8),
        Inches(0.8),
        "H大学图书馆馆藏结构优化中的取舍\n一般学科资源与民族地区文献建设如何兼顾",
        18,
        SUBTLE,
    )
    add_rect(slide, Inches(9.2), Inches(1.05), Inches(3.35), Inches(4.6), WHITE, MUTED, radius=True)
    add_textbox(slide, Inches(9.55), Inches(1.35), Inches(2.7), Inches(0.35), "这版 PPT 的表达更“像人在说话”", 13, ACCENT_WARM, True)
    add_paragraphs(
        slide,
        Inches(9.55),
        Inches(1.85),
        Inches(2.6),
        Inches(2.5),
        [
            ("我想把重点说得更直接一点：", 15, INK, True),
            ("这不是一份普通的馆藏统计。", 17, INK, True),
            ("它真正讨论的是，在有限预算里，学校到底要先保什么。", 17, INK, True),
        ],
    )
    add_rect(slide, Inches(0.82), Inches(5.7), Inches(8.0), Inches(0.8), ACCENT_LIGHT, ACCENT_LIGHT, radius=True)
    add_textbox(
        slide,
        Inches(1.0),
        Inches(5.95),
        Inches(7.5),
        Inches(0.3),
        "汇报口吻会比论文更鲜明，但仍保留教学案例需要的克制与判断感。",
        14,
        ACCENT,
        False,
    )
    add_footer(slide, 1, "版本：个人风格增强版")

    # Slide 2
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_slide_title(slide, "情境", "这不是一次普通的馆藏复盘", "先交代背景，再把真正的取舍抬出来")
    add_metric_card(slide, Inches(0.75), Inches(1.45), Inches(1.8), Inches(1.25), "142,847", "五年可识别新增馆藏", "2021—2025")
    add_metric_card(slide, Inches(2.75), Inches(1.45), Inches(1.8), Inches(1.25), "166,382", "原始记录总量", "含回溯 / 特藏 / 工具书")
    add_metric_card(slide, Inches(4.75), Inches(1.45), Inches(1.8), Inches(1.25), "4.98%", "民族文献相关新增占比", "7108 / 142847")
    add_rect(slide, Inches(7.1), Inches(1.38), Inches(5.45), Inches(4.65), WHITE, MUTED, radius=True)
    add_bullets(
        slide,
        Inches(7.35),
        Inches(1.75),
        Inches(4.9),
        Inches(3.9),
        [
            "我更想强调：这次讨论不是“书够不够”，而是“先为谁留位置”。",
            "学校既要服务一般学科教学科研，也承担民族地区人才培养、民族学科建设和文化传承任务。",
            "预算、空间、评价口径都有限，馆藏优化开始变成一道真正的选择题。",
        ],
        font_size=18,
    )
    add_rect(slide, Inches(0.8), Inches(3.0), Inches(5.8), Inches(2.9), WHITE, MUTED, radius=True)
    add_textbox(slide, Inches(1.05), Inches(3.25), Inches(5.2), Inches(0.3), "我会把这份案例理解成三个层次的矛盾", 15, ACCENT_WARM, True)
    add_paragraphs(
        slide,
        Inches(1.05),
        Inches(3.7),
        Inches(5.0),
        Inches(1.9),
        [
            ("1. 总量还在增长，但方向已经变了。", 17, INK, True),
            ("2. 特色资源不大，却最能代表学校是谁。", 17, INK, True),
            ("3. 如果继续用平均数看一切，很多价值会被看矮。", 17, INK, True),
        ],
    )
    add_footer(slide, 2)

    # Slide 3
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_slide_title(slide, "全馆变化", "五年里，馆藏没有停，重心却在悄悄换挡")
    add_column_chart(slide, Inches(0.78), Inches(1.45), Inches(7.0), Inches(4.8), years, "新增册数", new_vols, ACCENT)
    add_rect(slide, Inches(8.1), Inches(1.45), Inches(4.5), Inches(4.7), WHITE, MUTED, radius=True)
    add_bullets(
        slide,
        Inches(8.35),
        Inches(1.85),
        Inches(4.0),
        Inches(3.2),
        [
            "2022 年收缩最明显，2025 年重新扩张。",
            "即便波动，新增馆藏始终维持在 2.5 万册以上。",
            "我的感受是：规模没有停，但资源排序已经在换挡。",
        ],
        font_size=18,
    )
    add_metric_card(slide, Inches(8.35), Inches(4.9), Inches(1.9), Inches(0.95), "2022", "压力最大的一年", "同比 -22.94%")
    add_metric_card(slide, Inches(10.45), Inches(4.9), Inches(1.9), Inches(0.95), "2025", "明显修复的一年", "同比 +13.30%")
    add_footer(slide, 3)

    # Slide 4
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_slide_title(slide, "结构迁移", "从 2021 到 2025，资源配置更明显地向人文社科议题倾斜")
    add_bar_compare_chart(
        slide,
        Inches(0.78),
        Inches(1.45),
        Inches(7.35),
        Inches(4.95),
        compare_categories,
        {"2021占比%": compare_2021, "2025占比%": compare_2025},
    )
    add_rect(slide, Inches(8.4), Inches(1.45), Inches(4.15), Inches(4.95), WHITE, MUTED, radius=True)
    add_bullets(
        slide,
        Inches(8.65),
        Inches(1.85),
        Inches(3.65),
        Inches(2.9),
        [
            "B 类从 6.49% 升到 12.94%，是五年里最明显的上升项。",
            "T 类从 12.30% 降到 7.20%，F 类也明显回落。",
            "在我看来，这不是简单增减，而是学校任务在重排书单。",
        ],
        font_size=17,
    )
    add_rect(slide, Inches(8.65), Inches(4.95), Inches(3.65), Inches(0.95), ACCENT_LIGHT, ACCENT_LIGHT, radius=True)
    add_textbox(
        slide,
        Inches(8.85),
        Inches(5.18),
        Inches(3.2),
        Inches(0.45),
        "最值得追问的不是“哪个类多了”，而是“为什么它开始更重要”。",
        13,
        ACCENT,
        True,
    )
    add_footer(slide, 4)

    # Slide 5
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_slide_title(slide, "响应速度", "看上去是采购节奏，背后其实是资源响应速度")
    add_line_chart(slide, Inches(0.82), Inches(1.48), Inches(7.1), Inches(4.7), years, "一年内出版占比%", one_year_ratio, ACCENT_WARM)
    add_metric_card(slide, Inches(8.15), Inches(1.48), Inches(1.95), Inches(1.1), "1—2年", "出版时差中位数", "多数年份")
    add_metric_card(slide, Inches(10.3), Inches(1.48), Inches(1.95), Inches(1.1), "59.55%", "2025 一年内出版占比", "较 2023 明显回升")
    add_rect(slide, Inches(8.15), Inches(2.9), Inches(4.1), Inches(3.0), WHITE, MUTED, radius=True)
    add_bullets(
        slide,
        Inches(8.4),
        Inches(3.25),
        Inches(3.6),
        Inches(2.15),
        [
            "2023—2024 年，一年内出版占比分别降到 44.66% 和 48.22%。",
            "2025 年回升到 59.55%，说明响应速度在修复。",
            "如果我要判断服务温度，我会先看这组数据。",
        ],
        font_size=16,
    )
    add_footer(slide, 5)

    # Slide 6
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_slide_title(slide, "学校辨识度", "真正拉开学校辨识度的，不是总量，而是民族地区相关资源")
    add_metric_card(slide, Inches(0.78), Inches(1.48), Inches(1.8), Inches(1.2), "9015", "相关馆藏总记录", "L203 / L201 / L102 等")
    add_metric_card(slide, Inches(2.75), Inches(1.48), Inches(1.8), Inches(1.2), "7108", "近五年新增", "2021—2025")
    add_metric_card(slide, Inches(4.72), Inches(1.48), Inches(1.8), Inches(1.2), "4.98%", "占同期新增比重", "不大，但很关键")
    add_single_bar_chart(
        slide,
        Inches(0.8),
        Inches(3.0),
        Inches(6.85),
        Inches(3.0),
        location_names,
        location_counts,
        {0, 1, 4},
    )
    add_rect(slide, Inches(7.95), Inches(1.48), Inches(4.55), Inches(4.55), WHITE, MUTED, radius=True)
    add_bullets(
        slide,
        Inches(8.2),
        Inches(1.9),
        Inches(4.05),
        Inches(2.3),
        [
            "K 类历史地理是核心，同时覆盖 I 类文学、J 类艺术和 C 类社会科学总论。",
            "民族文献阅览与民大文库明显更偏专题研究型，价格也更高。",
            "它们规模不算最大，但最能代表这所学校是谁。",
        ],
        font_size=16,
    )
    add_metric_card(slide, Inches(8.2), Inches(4.5), Inches(1.9), Inches(1.0), "221.76", "民族文献阅览均价", "研究型更强")
    add_metric_card(slide, Inches(10.3), Inches(4.5), Inches(1.9), Inches(1.0), "262.04", "民大文库均价", "收藏属性更强")
    add_footer(slide, 6)

    # Slide 7
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_slide_title(slide, "使命承接", "这些资源之所以重要，是因为它们承接的是学校的特殊任务")
    card_w = Inches(2.9)
    card_h = Inches(1.65)
    positions = [
        (Inches(0.82), Inches(1.75)),
        (Inches(3.95), Inches(1.75)),
        (Inches(7.08), Inches(1.75)),
        (Inches(10.21), Inches(1.75)),
    ]
    card_texts = [
        ("面向民族地区人才培养", "不是只补文献，而是在补课程支撑。"),
        ("面向民族学科研究", "资源配置会直接影响研究能不能做深。"),
        ("面向文化传承与保护", "有些文献的价值，本来就不能只按均价来算。"),
        ("面向共同体意识教育", "专题资源已经在进入新的教学与研究场景。"),
    ]
    for (left, top), (title, body) in zip(positions, card_texts):
        add_rect(slide, left, top, card_w, card_h, WHITE, MUTED, radius=True)
        add_textbox(slide, left + Inches(0.18), top + Inches(0.18), card_w - Inches(0.36), Inches(0.38), title, 14, ACCENT, True)
        add_textbox(slide, left + Inches(0.18), top + Inches(0.65), card_w - Inches(0.36), Inches(0.7), body, 14, INK)
    add_rect(slide, Inches(1.0), Inches(4.3), Inches(11.2), Inches(1.15), ACCENT_LIGHT, ACCENT_LIGHT, radius=True)
    add_textbox(
        slide,
        Inches(1.25),
        Inches(4.66),
        Inches(10.7),
        Inches(0.42),
        "如果把这部分资源只算成“特色补充”，其实低估了它们的战略分量。",
        20,
        ACCENT_WARM,
        True,
        align=PP_ALIGN.CENTER,
    )
    add_bullets(
        slide,
        Inches(1.35),
        Inches(5.7),
        Inches(10.4),
        Inches(0.65),
        [
            "我会把它们理解成：学校办学定位在馆藏层面的可见证据。",
        ],
        font_size=15,
    )
    add_footer(slide, 7)

    # Slide 8
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_slide_title(slide, "现实难题", "难点不是看不见，而是必须先做选择")
    q_w = Inches(5.75)
    q_h = Inches(1.65)
    add_question_box(
        slide,
        Inches(0.82),
        Inches(1.7),
        q_w,
        q_h,
        "要不要单列预算？",
        "下一轮预算，是继续加码通用板块，还是把民族地区文献与特色馆藏单独拿出来保？",
    )
    add_question_box(
        slide,
        Inches(6.75),
        Inches(1.7),
        q_w,
        q_h,
        "要不要单列评价？",
        "民族地区相关资源，继续放在全馆平均数里，还是建立独立的年度评价口径？",
    )
    add_question_box(
        slide,
        Inches(0.82),
        Inches(3.7),
        q_w,
        q_h,
        "三处馆藏怎么分工？",
        "L203、L201、L102 分别更适合承担借阅、研究还是收藏任务？",
    )
    add_question_box(
        slide,
        Inches(6.75),
        Inches(3.7),
        q_w,
        q_h,
        "规则要不要分开？",
        "高价专题文献与课程复本，究竟该不该使用同一套论证和绩效标准？",
    )
    add_rect(slide, Inches(0.82), Inches(5.8), Inches(11.68), Inches(0.6), ACCENT, ACCENT, radius=True)
    add_textbox(
        slide,
        Inches(1.0),
        Inches(5.95),
        Inches(11.3),
        Inches(0.25),
        "教学型案例的价值，往往就藏在这些没有标准答案、但必须当场表态的问题里。",
        14,
        WHITE,
        True,
    )
    add_footer(slide, 8)

    # Slide 9
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_slide_title(slide, "我的判断", "先把口径分开，再把特色做实")
    step_texts = [
        "先分口径：常规新增、特色馆藏、工具书、捐赠资源分开统计。",
        "再单列：民族地区相关资源做年度看板，不再淹没在平均数里。",
        "再校准：把 B/C/D/T/F/G 的变化放回课程和学科任务里解释。",
        "再分治：高复本课程书与高价专题资料分别论证。",
        "最后用起来：让特色馆藏真正接到课程、导读、展示和数字化服务上。",
    ]
    top_positions = [
        (Inches(0.82), Inches(1.8)),
        (Inches(6.75), Inches(1.8)),
        (Inches(0.82), Inches(3.2)),
        (Inches(6.75), Inches(3.2)),
        (Inches(3.78), Inches(4.6)),
    ]
    card_w = Inches(5.45)
    card_h = Inches(1.1)
    for idx, (pos, text) in enumerate(zip(top_positions, step_texts), start=1):
        left, top = pos
        add_step_card(slide, left, top, card_w, card_h, str(idx), text)
    add_footer(slide, 9)

    # Slide 10
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_slide_title(slide, "课堂讨论", "如果你是资源建设部负责人，你会怎么选？")
    add_rect(slide, Inches(0.85), Inches(1.55), Inches(11.7), Inches(4.6), WHITE, MUTED, radius=True)
    questions = [
        "1. 在预算只小幅增长时，你会优先保哪一类资源？为什么？",
        "2. 民族地区文献该不该单列预算与评价口径？",
        "3. 你会怎样向馆长解释：为什么这 4.98% 不能只按平均数来看？",
    ]
    add_paragraphs(
        slide,
        Inches(1.2),
        Inches(2.0),
        Inches(10.8),
        Inches(2.5),
        [(q, 22, INK, True if idx == 0 else False) for idx, q in enumerate(questions)],
    )
    add_rect(slide, Inches(2.0), Inches(5.2), Inches(9.3), Inches(0.8), ACCENT_LIGHT, ACCENT_LIGHT, radius=True)
    add_textbox(
        slide,
        Inches(2.25),
        Inches(5.47),
        Inches(8.8),
        Inches(0.25),
        "我更愿意把这份案例留在“选择题”上，而不是急着把它做成“标准答案”。",
        18,
        ACCENT_WARM,
        True,
        align=PP_ALIGN.CENTER,
    )
    add_footer(slide, 10, "用途：课堂讨论 / 案例汇报 / 初步方案讨论")

    return prs


def main() -> None:
    prs = build_presentation()
    out = Path("馆藏结构优化_教学案例汇报_个人风格版.pptx")
    prs.save(str(out))
    print(out.name)
    print(len(prs.slides))


if __name__ == "__main__":
    main()
